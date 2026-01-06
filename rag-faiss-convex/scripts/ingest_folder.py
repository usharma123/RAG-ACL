import os
import sys
import json
import requests
from typing import Optional
from dotenv import load_dotenv
from openai import OpenAI

# Add parent dir to path for api module
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from api.faiss_store import FaissPerSourceStore

# Document parsing libraries (optional - graceful fallback)
try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    print("Note: PyPDF2 not installed, PDF files will be skipped")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Note: python-pptx not installed, PPTX files will be skipped")

try:
    from openpyxl import load_workbook
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False
    print("Note: openpyxl not installed, XLSX files will be skipped")

load_dotenv()

CONVEX_URL = os.environ["CONVEX_URL"].rstrip("/")
oa = OpenAI(
    api_key=os.environ["OPENAI_API_KEY"],
    base_url=os.getenv("OPENAI_BASE_URL"),
)
EMBED_MODEL = os.getenv("EMBED_MODEL", "text-embedding-3-small")

faiss_store = FaissPerSourceStore()

DATA_DIR = "data"
TENANT_ID = os.getenv("TENANT_ID", "acme")

def convex_mutation(path: str, args: dict):
    r = requests.post(
        f"{CONVEX_URL}/api/mutation",
        json={"path": path, "args": args, "format": "json"},
        headers={"Content-Type": "application/json"},
        timeout=60,
    )
    r.raise_for_status()
    data = r.json()
    if data.get("status") != "success":
        raise RuntimeError(data)
    return data["value"]

def chunk_text(text: str, max_chars: int = 1200, overlap_chars: int = 200):
    """
    Split text into chunks with overlap for better context preservation.
    
    Args:
        text: The text to chunk
        max_chars: Maximum characters per chunk (default 1200)
        overlap_chars: Characters to overlap between chunks (default 200)
    
    Returns:
        List of text chunks
    """
    # Split into paragraphs
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    
    if not paragraphs:
        return []
    
    chunks = []
    current_chunk = []
    current_len = 0
    
    for para in paragraphs:
        para_len = len(para)
        
        # If adding this paragraph exceeds max, save current chunk and start new one
        if current_chunk and current_len + para_len + 2 > max_chars:  # +2 for \n\n
            chunk_text = "\n\n".join(current_chunk)
            chunks.append(chunk_text)
            
            # Start new chunk with overlap from previous chunk
            # Take paragraphs from end of current_chunk that fit in overlap
            overlap_paras = []
            overlap_len = 0
            for p in reversed(current_chunk):
                if overlap_len + len(p) + 2 <= overlap_chars:
                    overlap_paras.insert(0, p)
                    overlap_len += len(p) + 2
                else:
                    break
            
            current_chunk = overlap_paras
            current_len = overlap_len
        
        current_chunk.append(para)
        current_len += para_len + 2
    
    # Don't forget the last chunk
    if current_chunk:
        chunks.append("\n\n".join(current_chunk))
    
    return chunks

def embed(texts):
    out = oa.embeddings.create(model=EMBED_MODEL, input=texts)
    return [d.embedding for d in out.data]

def parse_pdf(filepath: str) -> str:
    """Extract text from PDF file."""
    if not HAS_PDF:
        raise ValueError("PyPDF2 not installed")
    
    reader = PdfReader(filepath)
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text.strip())
    return "\n\n".join(text_parts)


def parse_pptx(filepath: str) -> str:
    """Extract text from PowerPoint file."""
    if not HAS_PPTX:
        raise ValueError("python-pptx not installed")
    
    prs = Presentation(filepath)
    text_parts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if len(slide_text) > 1:  # More than just the slide header
            text_parts.append("\n".join(slide_text))
    
    return "\n\n".join(text_parts)


def parse_xlsx(filepath: str) -> str:
    """Extract text from Excel file."""
    if not HAS_XLSX:
        raise ValueError("openpyxl not installed")
    
    wb = load_workbook(filepath, read_only=True, data_only=True)
    text_parts = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = [f"=== Sheet: {sheet_name} ==="]
        
        for row in sheet.iter_rows():
            row_values = []
            for cell in row:
                if cell.value is not None:
                    row_values.append(str(cell.value))
            if row_values:
                sheet_text.append(" | ".join(row_values))
        
        if len(sheet_text) > 1:
            text_parts.append("\n".join(sheet_text))
    
    wb.close()
    return "\n\n".join(text_parts)


def parse_file_to_text(filepath: str) -> Optional[tuple]:
    """
    Returns (title, text) or None if file type not supported.
    
    Supported formats:
    - .md/.txt: plain text
    - .json: Slack export format
    - .pdf: PDF documents
    - .pptx: PowerPoint presentations
    - .xlsx: Excel spreadsheets
    """
    name = os.path.basename(filepath)
    title = os.path.splitext(name)[0].replace("_", " ")
    ext = os.path.splitext(filepath)[1].lower()

    # Plain text files
    if ext in [".md", ".txt"]:
        with open(filepath, "r", encoding="utf-8") as f:
            return title, f.read()

    # Slack JSON export
    if ext == ".json":
        with open(filepath, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict) and "messages" in data:
            lines = [f"Channel: {data.get('channel','unknown')}"]
            for m in data["messages"]:
                lines.append(f"{m.get('user','?')}: {m.get('text','')}")
            return title, "\n".join(lines)
        return title, json.dumps(data)

    # PDF documents
    if ext == ".pdf":
        if not HAS_PDF:
            print(f"  Skipping {name} (PyPDF2 not installed)")
            return None
        try:
            text = parse_pdf(filepath)
            return title, text
        except Exception as e:
            print(f"  Error parsing {name}: {e}")
            return None

    # PowerPoint presentations
    if ext == ".pptx":
        if not HAS_PPTX:
            print(f"  Skipping {name} (python-pptx not installed)")
            return None
        try:
            text = parse_pptx(filepath)
            return title, text
        except Exception as e:
            print(f"  Error parsing {name}: {e}")
            return None

    # Excel spreadsheets
    if ext in [".xlsx", ".xls"]:
        if not HAS_XLSX:
            print(f"  Skipping {name} (openpyxl not installed)")
            return None
        try:
            text = parse_xlsx(filepath)
            return title, text
        except Exception as e:
            print(f"  Error parsing {name}: {e}")
            return None

    # Unknown file type - skip silently
    print(f"  Skipping {name} (unsupported file type: {ext})")
    return None

def ingest_doc(
    tenant_id: str,
    source_key: str,
    title: str,
    raw_text: str,
    source_url: str | None = None,
):
    doc_args = {
        "tenantId": tenant_id,
        "sourceKey": source_key,
        "title": title,
        "rawText": raw_text,
    }
    if source_url:
        doc_args["sourceUrl"] = source_url

    doc_id = convex_mutation("ingest:addDocument", doc_args)

    pieces = chunk_text(raw_text)
    vectors = embed(pieces)

    chunk_rows = [{"chunkIndex": i, "text": pieces[i]} for i in range(len(pieces))]
    chunk_ids = convex_mutation("ingest:addChunks", {
        "tenantId": tenant_id,
        "sourceKey": source_key,
        "docId": doc_id,
        "chunks": chunk_rows,
    })

    faiss_store.add(tenant_id, source_key, vectors, chunk_ids)
    print(f"Ingested [{source_key}] {title}  chunks={len(chunk_ids)}")

def main():
    if not os.path.isdir(DATA_DIR):
        raise RuntimeError(f"Missing {DATA_DIR}/ directory")

    total_docs = 0
    total_chunks = 0
    
    for source_key in os.listdir(DATA_DIR):
        src_dir = os.path.join(DATA_DIR, source_key)
        if not os.path.isdir(src_dir):
            continue

        print(f"\nProcessing source: {source_key}")
        
        for root, _, files in os.walk(src_dir):
            for fn in files:
                # Skip hidden files
                if fn.startswith("."):
                    continue
                    
                fp = os.path.join(root, fn)
                result = parse_file_to_text(fp)
                
                # Skip if file type not supported
                if result is None:
                    continue
                    
                title, text = result
                
                # Skip empty documents
                if not text.strip():
                    print(f"  Skipping {fn} (empty content)")
                    continue
                
                ingest_doc(TENANT_ID, source_key, title, text)
                total_docs += 1

    print(f"\n{'='*50}")
    print(f"Ingestion complete!")
    print(f"Total documents: {total_docs}")
    print(f"FAISS indexes stored in ./faiss_data/{TENANT_ID}/<source>.index")

if __name__ == "__main__":
    main()
