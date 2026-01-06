#!/usr/bin/env python3
"""
Generate sample documents for testing the RAG pipeline.
Includes: Markdown, JSON (Slack), PDF, PPTX, and XLSX files.

Install dependencies:
    pip install reportlab python-pptx openpyxl
"""
import os
import json
import random
from datetime import datetime, timedelta

BASE = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data")

# Try to import document generation libraries
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import inch
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False
    print("Note: reportlab not installed, skipping PDF generation")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Note: python-pptx not installed, skipping PPTX generation")

try:
    from openpyxl import Workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    print("Note: openpyxl not installed, skipping XLSX generation")


SOURCES = {
    "gdrive": [
        ("Q1_Budget_Notes.md", "Q1 Budget Notes"),
        ("Hiring_Plan.md", "Hiring Plan"),
    ],
    "confluence": [
        ("Onboarding_Runbook.md", "Onboarding Runbook"),
        ("Incident_Response.md", "Incident Response"),
    ],
    "slack": [
        ("eng_general.json", "Slack Export: #eng-general"),
        ("finance_team.json", "Slack Export: #finance-team"),
    ],
}


def ensure_dirs():
    os.makedirs(BASE, exist_ok=True)
    for src in ["gdrive", "confluence", "slack", "finance", "engineering", "hr"]:
        os.makedirs(os.path.join(BASE, src), exist_ok=True)


def write_md(path: str, title: str, body: str):
    with open(path, "w", encoding="utf-8") as f:
        f.write(f"# {title}\n\n{body}\n")


def make_paragraph(topic: str) -> str:
    bullets = [
        f"- Key point about {topic}: {random.choice(['process', 'policy', 'timeline', 'ownership'])}",
        f"- Risk: {random.choice(['latency', 'cost', 'compliance', 'handoff'])}",
        f"- Next step: {random.choice(['review', 'approve', 'ship', 'audit'])}",
    ]
    return "\n".join(bullets)


def generate_gdrive_docs():
    for filename, title in SOURCES["gdrive"]:
        body = "\n\n".join([
            "These are internal notes stored in Google Drive.",
            make_paragraph("budget"),
            make_paragraph("headcount"),
            "Confidential: Finance-related details included."
        ])
        write_md(os.path.join(BASE, "gdrive", filename), title, body)


def generate_confluence_docs():
    for filename, title in SOURCES["confluence"]:
        body = "\n\n".join([
            "This page is maintained in Confluence.",
            make_paragraph("onboarding"),
            make_paragraph("incident response"),
            "Reminder: follow the runbook exactly during incidents."
        ])
        write_md(os.path.join(BASE, "confluence", filename), title, body)


def generate_slack_exports():
    now = datetime.utcnow()
    channels = {
        "eng_general.json": ["deploy", "pager", "incident", "latency", "db migration"],
        "finance_team.json": ["reimbursements", "approval", "budget", "invoice", "expense policy"],
    }

    for filename, keywords in channels.items():
        messages = []
        for i in range(50):
            ts = (now - timedelta(minutes=15*i)).timestamp()
            messages.append({
                "user": random.choice(["alice", "bob", "carol", "dave"]),
                "ts": str(ts),
                "text": f"{random.choice(keywords)} update: {random.choice(['looks good', 'needs review', 'blocked', 'approved'])}"
            })

        path = os.path.join(BASE, "slack", filename)
        with open(path, "w", encoding="utf-8") as f:
            json.dump({"channel": filename.replace(".json",""), "messages": messages}, f, indent=2)


# ============ PDF Generation ============

def create_finance_pdf():
    """Create a sample Q4 Financial Report PDF."""
    if not HAS_REPORTLAB:
        return
    
    filepath = os.path.join(BASE, "finance", "Q4_Financial_Report_2024.pdf")
    c = canvas.Canvas(filepath, pagesize=letter)
    width, height = letter
    
    c.setFont("Helvetica-Bold", 24)
    c.drawString(1*inch, height - 1*inch, "Q4 2024 Financial Report")
    
    c.setFont("Helvetica", 14)
    c.drawString(1*inch, height - 1.5*inch, "Acme Corporation - Confidential")
    
    c.setFont("Helvetica", 11)
    y = height - 2.5*inch
    
    content = [
        "Executive Summary",
        "",
        "Q4 2024 was a strong quarter for Acme Corporation. Total revenue reached $42.5M,",
        "representing a 23% increase year-over-year. Key highlights include:",
        "",
        "Revenue: $42.5M (up 23% YoY)",
        "Gross Margin: 68% (up from 65% in Q3)",
        "Operating Expenses: $18.2M",
        "Net Income: $8.7M",
        "Cash Position: $125M",
        "",
        "Revenue Breakdown by Segment:",
        "",
        "1. Enterprise Software: $28.5M (67% of revenue)",
        "   - New enterprise deals: 12 contracts worth $15M ARR",
        "   - Expansion revenue: $8.5M from existing customers",
        "   - Churn rate: 2.1% (down from 3.2% in Q3)",
        "",
        "2. Professional Services: $9.2M (22% of revenue)",
        "   - Implementation projects: 45 completed",
        "   - Average project size: $205K",
        "   - Customer satisfaction: 4.8/5.0",
        "",
        "3. Support & Maintenance: $4.8M (11% of revenue)",
        "   - Support tickets resolved: 2,847",
        "   - Average resolution time: 4.2 hours",
        "   - SLA compliance: 99.2%",
        "",
        "Budget Allocation for Q1 2025:",
        "",
        "Engineering: $8.5M (40%)",
        "Sales & Marketing: $6.2M (29%)",
        "Operations: $3.8M (18%)",
        "G&A: $2.8M (13%)",
        "",
        "Key Risks and Mitigations:",
        "",
        "1. Market Competition: Increased investment in R&D",
        "2. Talent Retention: New equity refresh program announced",
        "3. Economic Uncertainty: Diversified customer base across industries",
        "",
        "Reimbursement Policy Update:",
        "All expense reports must be submitted within 30 days of the expense.",
        "Travel expenses require pre-approval for amounts over $500.",
        "The maximum daily meal allowance is $75 for domestic travel.",
    ]
    
    for line in content:
        if y < 1*inch:
            c.showPage()
            c.setFont("Helvetica", 11)
            y = height - 1*inch
        c.drawString(1*inch, y, line)
        y -= 14
    
    c.save()
    print(f"Created: {filepath}")


def create_gdrive_pdf():
    """Create a sample Product Roadmap PDF for gdrive source."""
    if not HAS_REPORTLAB:
        return
    
    filepath = os.path.join(BASE, "gdrive", "Product_Roadmap_2025.pdf")
    c = canvas.Canvas(filepath, pagesize=letter)
    width, height = letter
    
    c.setFont("Helvetica-Bold", 24)
    c.drawString(1*inch, height - 1*inch, "Product Roadmap 2025")
    
    c.setFont("Helvetica", 14)
    c.drawString(1*inch, height - 1.5*inch, "Acme Corporation - Product Team")
    
    c.setFont("Helvetica", 11)
    y = height - 2.5*inch
    
    content = [
        "Vision Statement",
        "",
        "Our mission is to build the most intuitive enterprise collaboration platform",
        "that unifies knowledge across all data sources with AI-powered intelligence.",
        "",
        "Q1 2025 Priorities:",
        "",
        "1. AI Search Enhancement (P0)",
        "   - Implement hybrid search (BM25 + vector)",
        "   - Add reranking with cross-encoder model",
        "   - Support for PDF, PPTX, XLSX extraction",
        "   - Target: 40% improvement in search relevance",
        "",
        "2. Real-time Collaboration (P0)",
        "   - Live cursors and presence indicators",
        "   - Collaborative document editing",
        "   - Comment threads with @mentions",
        "   - Target: Feature parity with Notion",
        "",
        "3. Enterprise Integrations (P1)",
        "   - Salesforce CRM connector",
        "   - Jira project sync",
        "   - Google Workspace deep integration",
        "   - Target: 5 new integrations by end of Q1",
        "",
        "Q2 2025 Priorities:",
        "",
        "1. Mobile App Launch (P0)",
        "   - iOS and Android native apps",
        "   - Offline mode support with sync",
        "   - Push notifications for mentions",
        "",
        "2. Advanced Analytics (P1)",
        "   - Usage dashboards for admins",
        "   - Knowledge gap analysis",
        "   - ROI reporting for enterprise customers",
        "",
        "Success Metrics:",
        "- DAU/MAU ratio > 60%",
        "- Search success rate > 85%",
        "- NPS > 50",
        "- Enterprise customer retention > 95%",
    ]
    
    for line in content:
        if y < 1*inch:
            c.showPage()
            c.setFont("Helvetica", 11)
            y = height - 1*inch
        c.drawString(1*inch, y, line)
        y -= 14
    
    c.save()
    print(f"Created: {filepath}")


# ============ PowerPoint Generation ============

def create_engineering_pptx():
    """Create a sample Engineering Architecture presentation."""
    if not HAS_PPTX:
        return
    
    filepath = os.path.join(BASE, "engineering", "System_Architecture_Overview.pptx")
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout
    
    # Slide 1: Title
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture Overview"
    p.font.size = Pt(44)
    p.font.bold = True
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Engineering Team - Q1 2025"
    p.font.size = Pt(24)
    
    # Slide 2: Architecture Components
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Architecture Components"
    p.font.size = Pt(32)
    p.font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    components = [
        ("Frontend Layer:", True),
        ("React 18 with TypeScript for type safety", False),
        ("Next.js 14 for server-side rendering and routing", False),
        ("TailwindCSS for utility-first styling", False),
        ("Deployed on Vercel Edge Network globally", False),
        ("", False),
        ("API Layer:", True),
        ("FastAPI Python 3.11 for high performance", False),
        ("GraphQL with Strawberry for flexible queries", False),
        ("Rate limiting: 1000 requests per minute per user", False),
        ("Authentication via JWT tokens with Convex Auth", False),
        ("", False),
        ("Data Layer:", True),
        ("Convex for real-time database and subscriptions", False),
        ("FAISS for vector similarity search", False),
        ("Redis for caching with 15 minute TTL", False),
        ("S3 for document and file storage", False),
    ]
    
    for i, (line, bold) in enumerate(components):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.bold = bold
    
    # Slide 3: Deployment Process
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Deployment Process"
    p.font.size = Pt(32)
    p.font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    deploy_info = [
        ("CI/CD Pipeline:", True),
        ("1. Developer pushes to feature branch", False),
        ("2. GitHub Actions runs unit, integration, and e2e tests", False),
        ("3. Code review required from 2 team members", False),
        ("4. Merge to main triggers automatic staging deploy", False),
        ("5. QA verification in staging environment (24hr soak test)", False),
        ("6. Production deploy via blue-green deployment strategy", False),
        ("", False),
        ("Rollback Procedure:", True),
        ("Automated rollback triggers if error rate exceeds 1%", False),
        ("Manual rollback available via CLI tool", False),
        ("Database migrations are always backwards-compatible", False),
        ("", False),
        ("Monitoring:", True),
        ("Datadog for APM, metrics, and centralized logs", False),
        ("PagerDuty for on-call alerts and escalation", False),
        ("SLO targets: 99.9% uptime, p99 latency under 200ms", False),
    ]
    
    for i, (line, bold) in enumerate(deploy_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.bold = bold
    
    # Slide 4: Tech Debt
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Tech Debt & Q1 Roadmap"
    p.font.size = Pt(32)
    p.font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    roadmap = [
        ("Current Tech Debt Priority:", True),
        ("1. Migrate legacy auth to Convex Auth - 2 weeks", False),
        ("2. Replace REST endpoints with GraphQL - 3 weeks", False),
        ("3. Upgrade React Query to TanStack Query v5 - 1 week", False),
        ("4. Implement proper error boundaries - 1 week", False),
        ("", False),
        ("Q1 2025 Engineering Initiatives:", True),
        ("Real-time collaboration features using WebSocket", False),
        ("Multi-region deployment to US-West and EU-Central", False),
        ("AI-powered search improvements with reranking", False),
        ("Mobile app MVP using React Native", False),
    ]
    
    for i, (line, bold) in enumerate(roadmap):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.bold = bold
    
    prs.save(filepath)
    print(f"Created: {filepath}")


def create_confluence_pptx():
    """Create a sample Incident Response presentation."""
    if not HAS_PPTX:
        return
    
    filepath = os.path.join(BASE, "confluence", "Incident_Response_Playbook.pptx")
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    
    # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Incident Response Playbook"
    p.font.size = Pt(44)
    p.font.bold = True
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "On-Call Engineering Guide"
    p.font.size = Pt(24)
    
    # Severity Levels
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Incident Severity Levels"
    p.font.size = Pt(32)
    p.font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    severity_info = [
        ("SEV1 - Critical (Response time: 15 minutes)", True),
        ("Complete service outage affecting all users", False),
        ("Data breach or security incident", False),
        ("Revenue-impacting issue over $10K per hour", False),
        ("", False),
        ("SEV2 - High (Response time: 30 minutes)", True),
        ("Partial service degradation affecting over 25% of users", False),
        ("Critical feature completely unavailable", False),
        ("Performance degradation more than 5x normal latency", False),
        ("", False),
        ("SEV3 - Medium (Response time: 4 hours)", True),
        ("Minor feature issues with workarounds available", False),
        ("Affecting less than 10% of users", False),
        ("", False),
        ("SEV4 - Low (Response time: 24 hours)", True),
        ("Cosmetic issues and documentation errors", False),
        ("Minor bugs with easy workarounds", False),
    ]
    
    for i, (line, bold) in enumerate(severity_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.bold = bold
    
    # Response Process
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Incident Response Process"
    p.font.size = Pt(32)
    p.font.bold = True
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    
    process = [
        ("Step 1: DETECT", True),
        ("PagerDuty alert received on phone", False),
        ("Acknowledge alert within SLA timeframe", False),
        ("Join #incident-response Slack channel immediately", False),
        ("", False),
        ("Step 2: ASSESS", True),
        ("Check Datadog dashboards for metrics", False),
        ("Identify which services are affected", False),
        ("Determine severity level based on impact", False),
        ("Page additional responders if needed", False),
        ("", False),
        ("Step 3: MITIGATE", True),
        ("Apply immediate fix or rollback deployment", False),
        ("Communicate status updates to stakeholders", False),
        ("Update public status page if customer-facing", False),
        ("", False),
        ("Step 4: RESOLVE AND REVIEW", True),
        ("Verify the fix is working correctly", False),
        ("Monitor metrics for 30 minutes post-fix", False),
        ("Schedule post-mortem within 48 hours", False),
    ]
    
    for i, (line, bold) in enumerate(process):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.bold = bold
    
    prs.save(filepath)
    print(f"Created: {filepath}")


# ============ Excel Generation ============

def create_hr_xlsx():
    """Create a sample HR Employee Directory Excel file."""
    if not HAS_OPENPYXL:
        return
    
    filepath = os.path.join(BASE, "hr", "Employee_Directory_2025.xlsx")
    wb = Workbook()
    
    # Sheet 1: Employee Directory
    ws = wb.active
    ws.title = "Employee Directory"
    
    headers = ["Employee ID", "Name", "Department", "Title", "Location", "Start Date", "Manager", "Email"]
    ws.append(headers)
    
    employees = [
        ["E001", "Sarah Chen", "Engineering", "VP of Engineering", "San Francisco", "2020-03-15", "CEO", "sarah.chen@acme.com"],
        ["E002", "James Wilson", "Engineering", "Senior Software Engineer", "San Francisco", "2021-06-01", "Sarah Chen", "james.wilson@acme.com"],
        ["E003", "Maria Garcia", "Engineering", "Software Engineer", "Remote", "2022-01-10", "Sarah Chen", "maria.garcia@acme.com"],
        ["E004", "David Kim", "Engineering", "DevOps Engineer", "San Francisco", "2021-09-20", "Sarah Chen", "david.kim@acme.com"],
        ["E005", "Emily Brown", "Finance", "CFO", "New York", "2019-11-01", "CEO", "emily.brown@acme.com"],
        ["E006", "Michael Johnson", "Finance", "Financial Analyst", "New York", "2022-03-15", "Emily Brown", "michael.johnson@acme.com"],
        ["E007", "Lisa Wang", "HR", "HR Director", "San Francisco", "2020-08-01", "CEO", "lisa.wang@acme.com"],
        ["E008", "Robert Taylor", "HR", "Recruiter", "Remote", "2023-02-01", "Lisa Wang", "robert.taylor@acme.com"],
        ["E009", "Jennifer Martinez", "Sales", "VP of Sales", "New York", "2020-05-10", "CEO", "jennifer.martinez@acme.com"],
        ["E010", "Chris Anderson", "Sales", "Account Executive", "Chicago", "2022-07-01", "Jennifer Martinez", "chris.anderson@acme.com"],
    ]
    
    for emp in employees:
        ws.append(emp)
    
    # Sheet 2: Benefits Summary
    ws2 = wb.create_sheet("Benefits Summary")
    ws2.append(["Benefit Type", "Description", "Eligibility", "Company Contribution"])
    
    benefits = [
        ["Health Insurance", "PPO and HMO options through Blue Cross Blue Shield", "All full-time employees from day one", "Company pays 80% of premium"],
        ["Dental Insurance", "Delta Dental PPO plan", "All full-time employees from day one", "Company pays 100% of premium"],
        ["Vision Insurance", "VSP Vision Care plan", "All full-time employees from day one", "Company pays 100% of premium"],
        ["401(k) Retirement", "Fidelity retirement plan with immediate vesting", "Eligible after 90 days of employment", "Company matches 4% of salary"],
        ["Life Insurance", "Coverage equal to 2x annual salary", "All full-time employees", "Company pays 100% of premium"],
        ["PTO Policy", "Unlimited PTO with minimum 15 days recommended", "All employees", "Not applicable"],
        ["Parental Leave", "16 weeks paid leave for all new parents", "After 1 year of employment", "100% of salary paid"],
        ["Learning Budget", "Annual professional development stipend", "All full-time employees", "$2,000 per year"],
        ["Home Office", "One-time equipment allowance for remote setup", "Remote employees only", "$1,500 one-time"],
        ["Commuter Benefits", "Pre-tax transit and parking benefits", "Office-based employees", "Up to $300 per month"],
    ]
    
    for benefit in benefits:
        ws2.append(benefit)
    
    # Sheet 3: Onboarding Checklist
    ws3 = wb.create_sheet("Onboarding Checklist")
    ws3.append(["Timeline", "Task", "Owner", "Notes"])
    
    onboarding = [
        ["Day 1 Morning", "Complete I-9 verification and tax forms", "HR Team", "Bring two forms of identification"],
        ["Day 1 Morning", "Receive laptop and equipment from IT", "IT Team", "MacBook Pro with monitor and accessories"],
        ["Day 1 Morning", "Set up email and Slack accounts", "IT Team", "Use your @acme.com email address"],
        ["Day 1 Afternoon", "Meet with your direct manager", "Your Manager", "30-minute introduction meeting"],
        ["Day 1 Afternoon", "Team welcome lunch", "Your Manager", "Expensed up to $30 per person"],
        ["Week 1", "Complete mandatory security training", "HR Team", "Required for all new employees"],
        ["Week 1", "Set up local development environment", "Engineering", "Follow the wiki setup guide"],
        ["Week 1", "Review employee handbook in Confluence", "HR Team", "Acknowledge receipt in Workday"],
        ["Week 1", "Schedule 1:1 with skip-level manager", "Your Manager", "15-minute introductory meeting"],
        ["Week 2", "Complete product training sessions", "Product Team", "2-hour interactive session"],
        ["Week 2", "Shadow a customer call with sales", "Sales Team", "Learn about customer needs"],
        ["Month 1", "30-day check-in with HR", "HR Team", "Complete feedback survey"],
        ["Month 2", "Complete your first assigned project", "Your Manager", "Scope defined in onboarding document"],
        ["Month 3", "90-day performance review", "Your Manager", "Formal performance discussion"],
    ]
    
    for task in onboarding:
        ws3.append(task)
    
    # Sheet 4: Office Hours
    ws4 = wb.create_sheet("Office Information")
    ws4.append(["Office Location", "Address", "Hours", "Contact", "Amenities"])
    
    offices = [
        ["San Francisco HQ", "123 Market Street, SF, CA 94105", "Mon-Fri 8am-6pm", "sf-office@acme.com", "Gym, cafeteria, rooftop deck, bike storage"],
        ["New York Office", "456 Broadway, New York, NY 10013", "Mon-Fri 9am-7pm", "ny-office@acme.com", "Gym, coffee bar, phone booths"],
        ["Chicago Office", "789 Michigan Ave, Chicago, IL 60611", "Mon-Fri 8am-5pm", "chi-office@acme.com", "Coffee bar, collaboration spaces"],
    ]
    
    for office in offices:
        ws4.append(office)
    
    wb.save(filepath)
    print(f"Created: {filepath}")


def main():
    print("Generating sample documents for RAG testing...\n")
    
    ensure_dirs()
    
    # Original markdown and JSON docs
    generate_gdrive_docs()
    generate_confluence_docs()
    generate_slack_exports()
    
    # New PDF, PPTX, XLSX docs
    create_finance_pdf()
    create_gdrive_pdf()
    create_engineering_pptx()
    create_confluence_pptx()
    create_hr_xlsx()
    
    print("\nDone! Sample documents created in data/ directory.")
    print("\nTo generate PDF/PPTX/XLSX files, install:")
    print("  pip install reportlab python-pptx openpyxl")
    print("\nThen run 'python scripts/ingest_folder.py' to ingest into the RAG system.")


if __name__ == "__main__":
    main()
